import json
import os
import subprocess
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def run_experiments(config_file):
    # 读取配置文件
    with open(config_file) as f:
        configs = json.load(f)

    # 创建结果目录
    os.makedirs("results", exist_ok=True)

    # 运行每个配置的可执行文件
    for idx, config in enumerate(configs):
        output_dir = f"results/config_{idx}"
        os.makedirs(output_dir, exist_ok=True)
        
        # 构建命令行参数（根据实际可执行文件参数调整）
        cmd = ["./your_executable"]  # 替换为实际可执行文件路径
        cmd += [f"--{k}={v}" for k, v in config.items()]
        
        # 运行程序
        try:
            subprocess.run(cmd, cwd=output_dir, check=True)
        except subprocess.CalledProcessError as e:
            print(f"Error running config {idx}: {e}")

def create_report():
    prs = Presentation()
    plot_paths = []

    # 遍历所有结果目录
    for idx, config_dir in enumerate(sorted(os.listdir("results"))):
        csv_path = os.path.join("results", config_dir, "output.csv")  # 替换为实际CSV文件名
        
        if not os.path.exists(csv_path):
            continue

        # 读取CSV数据
        df = pd.read_csv(csv_path, header=None)
        
        # 提取第一行和最后三行
        selected_data = pd.concat([df.head(1), df.tail(3)])
        
        # 绘制图表
        plt.figure(figsize=(8, 6))
        plt.plot(selected_data.index.values, selected_data.iloc[:, 0],  # 选择要绘制的列
                 marker='o', linestyle='-', linewidth=2)
        plt.title(f"Configuration {idx}")
        plt.xlabel("Data Points")
        plt.ylabel("Values")
        
        # 保存图片
        plot_path = f"plot_{idx}.png"
        plt.savefig(plot_path, dpi=300, bbox_inches='tight')
        plt.close()
        plot_paths.append(plot_path)

    # 将图表添加到PPT
    for path in plot_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白版式
        slide.shapes.add_picture(
            path,
            left=Inches(1),
            top=Inches(1),
            width=Inches(8),
            height=Inches(6)
        )
        os.remove(path)  # 清理临时图片

    prs.save("experiment_report.pptx")

if __name__ == "__main__":
    # 第一步：运行所有实验
    run_experiments("config.json")
    
    # 第二步：生成报告
    create_report()
