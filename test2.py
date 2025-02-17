import os
import argparse
import subprocess
import time
import csv
from threading import Thread
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def run_process(exe_path, folder):
    input_file = os.path.join(folder, "input.txt")
    output_file = os.path.join(folder, "output.txt")
    
    # 读取输入参数
    with open(input_file) as f:
        params = [line.strip() for line in f.readlines()]

    # 启动进程并重定向输出
    with open(output_file, "w") as out_f:
        proc = subprocess.Popen(
            [exe_path],
            stdin=subprocess.PIPE,
            stdout=out_f,
            stderr=subprocess.STDOUT,
            text=True,
            cwd=folder
        )

        # 每秒发送一个参数的线程
        def send_params():
            for param in params:
                time.sleep(1)
                proc.stdin.write(param + "\n")
                proc.stdin.flush()
            proc.stdin.close()

        sender = Thread(target=send_params)
        sender.start()

        proc.wait()
        sender.join()

def run_experiments(exe_path, folders):
    processes = []
    for folder in folders:
        if not os.path.isdir(folder):
            print(f"警告：跳过无效目录 {folder}")
            continue
        
        # 为每个文件夹启动独立进程
        p = Thread(target=run_process, args=(exe_path, folder))
        p.start()
        processes.append(p)

    # 等待所有进程完成
    for p in processes:
        p.join()

def create_report(template_path, folders):
    prs = Presentation(template_path)
    
    for folder in folders:
        csv_path = os.path.join(folder, "data.csv")  # 根据实际情况修改CSV文件名
        if not os.path.exists(csv_path):
            continue

        # 读取CSV数据
        with open(csv_path) as f:
            reader = csv.reader(f)
            rows = list(reader)

        # 提取数据：第一行时间，最后三行角度
        if len(rows) < 4:
            continue
            
        time_data = [float(x) for x in rows[0]]
        angle_data = [
            [float(x) for x in rows[-3]],
            [float(x) for x in rows[-2]],
            [float(x) for x in rows[-1]]
        ]

        # 创建散点图
        plt.figure(figsize=(10, 6))
        for i, angles in enumerate(angle_data):
            plt.scatter(time_data, angles, label=f'Run {i+1}')
            
        plt.title(os.path.basename(folder))
        plt.xlabel("Time")
        plt.ylabel("Angle")
        plt.legend()
        plt.grid(True)

        # 保存临时图片
        plot_path = os.path.join(folder, "temp_plot.png")
        plt.savefig(plot_path, dpi=300, bbox_inches='tight')
        plt.close()

        # 添加图片到PPT
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 根据模板调整版式
        slide.shapes.add_picture(
            plot_path,
            left=Inches(1),
            top=Inches(2),
            width=Inches(8),
            height=Inches(4.5)
        )
        os.remove(plot_path)

    prs.save("final_report.pptx")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("exe_path", help="Path to the executable file")
    parser.add_argument("folders", nargs="+", help="List of working directories")
    parser.add_argument("--template", default="template.pptx", help="PPT template path")
    args = parser.parse_args()

    # 运行所有实验
    run_experiments(args.exe_path, args.folders)
    
    # 生成报告
    create_report(args.template, args.folders)
